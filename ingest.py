"""
ingest.py - Fixed:
  - chunks now carry 'topic' and 'chunk_id' fields (needed by quiz router)
  - embed_texts batches all texts at once (much faster)
  - retrieve_chunks returns score in each chunk dict
  - get_all_topics() and get_chunks_by_topic() helpers added
"""

import os, pickle, re, hashlib
from typing import List
import fitz, faiss, numpy as np
try:
    from pptx import Presentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False
from sentence_transformers import SentenceTransformer
from database import SessionLocal, EmbeddingCache, Document
from datetime import datetime

CHUNK_SIZE   = 512
CHUNK_OVERLAP = 64
FAISS_PATH   = "vector_store/faiss.index"
CHUNKS_PATH  = "vector_store/chunks.pkl"
os.makedirs("vector_store", exist_ok=True)

model = SentenceTransformer("all-MiniLM-L6-v2")
index = None
chunks_store: List[dict] = []


# ── store helpers ──────────────────────────────────────────────────────────────
def load_store():
    global index, chunks_store
    if os.path.exists(FAISS_PATH) and os.path.exists(CHUNKS_PATH):
        index = faiss.read_index(FAISS_PATH)
        with open(CHUNKS_PATH, "rb") as f:
            chunks_store = pickle.load(f)
    else:
        index = faiss.IndexFlatIP(384)
        chunks_store = []

def save_store():
    faiss.write_index(index, FAISS_PATH)
    with open(CHUNKS_PATH, "wb") as f:
        pickle.dump(chunks_store, f)


# ── Text extraction (PDF + PPTX) ───────────────────────────────────────────────
def _extract_pdf(path: str) -> str:
    doc  = fitz.open(path)
    text = "\n".join(p.get_text() for p in doc)
    doc.close()
    return text

def _extract_pptx(path: str) -> str:
    if not _PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")
    prs, slides = Presentation(path), []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(slide.shapes.title.text.strip())
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
        slides.append("\n".join(parts))
    return "\n\n".join(slides)

def extract_text(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    return _extract_pptx(path) if ext in (".pptx", ".ppt") else _extract_pdf(path)


# ── Chunking (sentence-aware sliding window) ───────────────────────────────────
def _infer_topic(text: str) -> str:
    first = text.strip().split("\n")[0][:80]
    clean = re.sub(r'[^a-zA-Z0-9 ]', '', first).strip()
    return clean if len(clean) > 3 else "General"

def split_chunks(text: str, source: str) -> List[dict]:
    text = re.sub(r'\s+', ' ', text).strip()
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks, current, length, cid = [], [], 0, 0
    for s in sentences:
        sl = len(s.split())
        if length + sl > CHUNK_SIZE and current:
            chunk_text = " ".join(current)
            chunks.append({"text": chunk_text, "source": source,
                            "chunk_id": cid, "topic": _infer_topic(chunk_text)})
            cid += 1
            overlap = " ".join(chunk_text.split()[-CHUNK_OVERLAP:])
            current = overlap.split(); length = len(current)
        current.append(s); length += sl
    if current:
        chunk_text = " ".join(current)
        chunks.append({"text": chunk_text, "source": source,
                        "chunk_id": cid, "topic": _infer_topic(chunk_text)})
    return chunks


# ── Embedding with DB cache (batch) ────────────────────────────────────────────
def embed_texts(texts: List[str]) -> np.ndarray:
    db = SessionLocal()
    result = [None] * len(texts)
    to_compute_idx, to_compute_texts = [], []

    for i, t in enumerate(texts):
        h = hashlib.md5(t.encode()).hexdigest()
        cached = db.query(EmbeddingCache).filter_by(text_hash=h).first()
        if cached:
            result[i] = np.array(cached.embedding, dtype=np.float32)
        else:
            to_compute_idx.append(i)
            to_compute_texts.append(t)

    if to_compute_texts:
        embs = model.encode(to_compute_texts, batch_size=64,
                            normalize_embeddings=True, show_progress_bar=False)
        for idx, text, emb in zip(to_compute_idx, to_compute_texts, embs):
            h = hashlib.md5(text.encode()).hexdigest()
            db.add(EmbeddingCache(text_hash=h, embedding=emb.tolist()))
            result[idx] = emb.astype(np.float32)
        db.commit()

    db.close()
    return np.array(result, dtype=np.float32)


# ── Main ingestion pipeline ────────────────────────────────────────────────────
def ingest_pdf(path: str, filename: str) -> dict:
    load_store()
    text = extract_text(path)
    if not text.strip():
        raise ValueError("No text extracted from file.")

    chunks = split_chunks(text, filename)
    embs = embed_texts([c["text"] for c in chunks])
    index.add(embs)
    chunks_store.extend(chunks)
    save_store()

    topics = list(set(c["topic"] for c in chunks))

    # DB record
    db = SessionLocal()
    existing = db.query(Document).filter_by(filename=filename).first()
    if existing:
        existing.total_chunks = len(chunks); existing.topics = topics
    else:
        db.add(Document(filename=filename, total_chunks=len(chunks), topics=topics))
    db.commit(); db.close()

    return {"status": "success", "chunks": chunks, "total_chunks": len(chunks), "topics": topics}


# ── Retrieval ──────────────────────────────────────────────────────────────────
def retrieve_chunks(query: str, top_k: int = 5) -> List[dict]:
    load_store()
    if index is None or index.ntotal == 0:
        return []
    q_emb = model.encode([query], normalize_embeddings=True).astype(np.float32)
    scores, indices = index.search(q_emb, min(top_k, index.ntotal))
    results = []
    for score, i in zip(scores[0], indices[0]):
        if 0 <= i < len(chunks_store):
            c = dict(chunks_store[i]); c["score"] = float(score)
            results.append(c)
    return results


# ── Topic helpers ──────────────────────────────────────────────────────────────
def get_all_topics() -> List[str]:
    load_store()
    return list(set(c.get("topic", "General") for c in chunks_store))

def get_chunks_by_topic(topic: str) -> List[dict]:
    load_store()
    return [c for c in chunks_store if c.get("topic") == topic]